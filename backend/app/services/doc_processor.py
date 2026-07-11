import os
os.environ["PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"] = "True"
import fitz  # PyMuPDF
from pathlib import Path
from typing import List, Dict, Any, Tuple
from ..utils.logger import get_logger
from ..utils.helpers import clean_text

logger = get_logger("DocProcessor")

# Attempt importing PaddleOCR with soft fallback
try:
    from paddleocr import PaddleOCR
    logger.info("PaddleOCR library imported successfully.")
    # Initialize engine lazily to avoid heavy start lockups
    ocr_engine = None
except ImportError:
    logger.warning("PaddleOCR library not installed. OCR falls back to basic metadata scanning.")
    ocr_engine = "FAILED"

class DocumentProcessor:
    @staticmethod
    def initialize_ocr():
        """Lazily initialize the OCR engine to speed up server startups."""
        global ocr_engine
        if ocr_engine == "FAILED":
            return None
        if ocr_engine is not None:
            return ocr_engine
        try:
            # ch_en is the default english/chinese model, we can switch to English
            ocr_engine = PaddleOCR(use_angle_cls=True, lang='en')
            logger.info("PaddleOCR engine loaded offline model parameters.")
            return ocr_engine
        except Exception as e:
            logger.error(f"Failed loading PaddleOCR model: {str(e)}")
            ocr_engine = "FAILED"  # Cache failure so we don't block subsequent requests
            return None

    @classmethod
    def extract_text_from_pdf(cls, file_path: Path) -> List[Tuple[int, str]]:
        """Extracts text page-by-page. Returns list of tuples (page_num, text)."""
        pages_content = []
        try:
            doc = fitz.open(str(file_path))
            for i in range(len(doc)):
                page = doc[i]
                text = page.get_text()
                
                # Check if page is scanned/empty text
                if not text.strip():
                    # Attempt OCR on page if PaddleOCR is available
                    ocr = cls.initialize_ocr()
                    if ocr:
                        logger.info(f"Page {i+1} has no selectable text, running PaddleOCR...")
                        pix = page.get_pixmap()
                        # Save pixmap temp
                        temp_img = file_path.parent / f"_temp_page_{i+1}.png"
                        pix.save(str(temp_img))
                        
                        result = ocr.ocr(str(temp_img), cls=True)
                        temp_img.unlink(missing_ok=True)
                        
                        ocr_texts = []
                        if result and result[0]:
                            for line in result[0]:
                                ocr_texts.append(line[1][0]) # Extract text tokens
                        text = " ".join(ocr_texts)
                        logger.info(f"OCR successfully extracted {len(text)} characters on page {i+1}.")
                
                pages_content.append((i + 1, clean_text(text)))
            doc.close()
        except Exception as e:
            logger.error(f"Failed parsing PDF {file_path.name}: {str(e)}")
        return pages_content

    @classmethod
    def extract_text_from_docx(cls, file_path: Path) -> str:
        """Parses Microsoft Word docx file text contents."""
        try:
            import docx
            doc = docx.Document(str(file_path))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return clean_text("\n".join(full_text))
        except ImportError:
            logger.error("python-docx is not installed.")
            return "Word parsing error (library missing)"
        except Exception as e:
            logger.error(f"Failed parsing Word Doc {file_path.name}: {str(e)}")
            return ""

    @classmethod
    def extract_text_from_pptx(cls, file_path: Path) -> str:
        """Parses Microsoft PowerPoint slides pptx contents."""
        try:
            import pptx
            prs = pptx.Presentation(str(file_path))
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return clean_text("\n".join(full_text))
        except ImportError:
            logger.error("python-pptx is not installed.")
            return "Powerpoint parsing error (library missing)"
        except Exception as e:
            logger.error(f"Failed parsing PowerPoint {file_path.name}: {str(e)}")
            return ""

    @classmethod
    def extract_text_from_image(cls, file_path: Path) -> str:
        """Runs PaddleOCR directly on standard image files."""
        ocr = cls.initialize_ocr()
        if not ocr:
            logger.warning("OCR engine not available. Image text cannot be read offline.")
            return "Image text extraction skipped (OCR unavailable)"
        try:
            logger.info(f"Running OCR on image {file_path.name}")
            result = ocr.ocr(str(file_path), cls=True)
            ocr_texts = []
            if result and result[0]:
                for line in result[0]:
                    ocr_texts.append(line[1][0])
            return clean_text(" ".join(ocr_texts))
        except Exception as e:
            logger.error(f"Failed OCR on image {file_path.name}: {str(e)}")
            return ""

    @classmethod
    def extract_text_from_txt(cls, file_path: Path) -> str:
        """Extracts text from raw TXT/MD offline document."""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return clean_text(f.read())
        except Exception as e:
            logger.error(f"Failed parsing plain text file {file_path.name}: {str(e)}")
            return ""

    @classmethod
    def process_file(cls, file_path: Path, file_type: str) -> List[Tuple[int, str]]:
        """Processes file based on type, returns a list of (page_num, text_content)."""
        logger.info(f"Indexing file {file_path.name} of type {file_type}")
        
        if file_type == "pdf":
            return cls.extract_text_from_pdf(file_path)
        elif file_type == "docx":
            return [(1, cls.extract_text_from_docx(file_path))]
        elif file_type == "ppt":
            return [(1, cls.extract_text_from_pptx(file_path))]
        elif file_type in ["image", "png", "jpg", "jpeg"]:
            return [(1, cls.extract_text_from_image(file_path))]
        elif file_type in ["txt", "md", "json"]:
            return [(1, cls.extract_text_from_txt(file_path))]
        else:
            logger.warning(f"Unsupported file format index query: {file_type}")
            return []

    @staticmethod
    def auto_detect_subject(text: str) -> List[str]:
        """Runs simple heuristic text parser to auto generate tags for the file."""
        if not text:
            return ["General"]
        
        tags = []
        keywords = {
            "Transformer": ["attention", "transformer", "bert", "gpt", "sequence"],
            "Finance": ["revenue", "projected", "finance", "forecast", "budget", "quarterly"],
            "Security": ["offline", "isolated", "encryption", "firewall", "credentials"],
            "System Architecture": ["infrastructure", "vector database", "cache", "pipelines", "rag"],
            "OCR Parser": ["scanned", "handwritten", "tesseract", "image", "capture"]
        }
        
        text_lower = text.lower()
        for tag, words in keywords.items():
            for word in words:
                if word in text_lower:
                    tags.append(tag)
                    break
        
        if not tags:
            tags.append("General")
            
        return tags[:3] # Limit to top 3 tags
